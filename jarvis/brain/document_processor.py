#!/usr/bin/env python3
"""
DocumentProcessor — processes files dropped into ACCOUNTS/{name}/DOCUMENTS/.

When you drop an RFP, contract, NDA, brief, or any document into an account's
DOCUMENTS/ folder, JARVIS reads it with NVIDIA, extracts intelligence,
and updates the account's context files automatically.

Supported formats: .pdf .docx .txt .md
"""

import asyncio
import json
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional

from jarvis.utils.logger import JARVISLogger
from jarvis.utils.event_bus import Event, EventBus
from jarvis.utils.config import ConfigManager


SUPPORTED_EXTS = {".pdf", ".docx", ".txt", ".md", ".doc", ".pptx", ".xlsx"}

# Audio/video formats handled by dedicated NVIDIA models
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".ogg", ".flac"}
VIDEO_EXTS = {".mp4", ".mov", ".avi", ".mkv", ".webm"}

# Char threshold above which we switch to the 1M-token Nemotron 120B model
LONG_CONTENT_THRESHOLD = 8000

# Which sections downstream should refresh when a file is dropped in each section
SECTION_CASCADE_MAP = {
    "DISCOVERY":        ["fill_discovery", "fill_demo_strategy", "fill_value_architecture"],
    "BATTLECARD":       ["fill_battlecard"],
    "DEMO_STRATEGY":    ["fill_demo_strategy"],
    "RISK_REPORT":      ["fill_risk_report"],
    "NEXT_STEPS":       ["fill_next_steps"],
    "VALUE_ARCHITECTURE": ["fill_value_architecture"],
}

DOCUMENT_TYPE_PROMPTS = {
    "rfp": (
        "This is an RFP (Request for Proposal). Extract: "
        "1) key requirements and evaluation criteria "
        "2) decision timeline and submission deadline "
        "3) budget range if mentioned "
        "4) competitors likely being evaluated "
        "5) top 5 win themes we should address "
        "6) MEDDPICC signals (decision process, paper process, metrics). "
        "Return as JSON."
    ),
    "contract": (
        "This is a contract or agreement. Extract: "
        "1) deal value and payment terms "
        "2) key obligations and milestones "
        "3) renewal and cancellation terms "
        "4) stakeholders and signatories. "
        "Return as JSON."
    ),
    "brief": (
        "This is a company or project brief. Extract: "
        "1) company background and size "
        "2) stated pain points or goals "
        "3) key people mentioned "
        "4) budget or investment signals "
        "5) competitive mentions. "
        "Return as JSON."
    ),
    "email": (
        "This is an email thread. Extract: "
        "1) all people mentioned with roles "
        "2) action items with owners "
        "3) deal status signals "
        "4) MEDDPICC signals "
        "5) next steps agreed. "
        "Return as JSON."
    ),
    "general": (
        "Extract from this document: "
        "1) company names and people mentioned "
        "2) key facts relevant to a sales deal "
        "3) any action items or commitments "
        "4) budget or timeline signals "
        "5) competitive mentions. "
        "Return as JSON."
    ),
}


class DocumentProcessor:
    """Processes files dropped into account DOCUMENTS/ or EMAILS/ folders."""

    def __init__(self, config_manager: ConfigManager, event_bus: EventBus):
        self.config = config_manager
        self.event_bus = event_bus
        self.logger = JARVISLogger("brain.documents")
        self._llm_client = None

    async def start(self):
        self.event_bus.subscribe("document.added",   self._on_document_added)
        self.event_bus.subscribe("email.added",       self._on_email_added)
        self.event_bus.subscribe("rfi.file.added",    self._on_rfi_file_added)
        self.event_bus.subscribe("presales.file.added", self._on_presales_file_added)
        self.logger.info("DocumentProcessor started")

    async def stop(self):
        self.logger.info("DocumentProcessor stopped")

    async def _on_document_added(self, event: Event):
        await self.process_document(
            file_path=Path(event.data["path"]),
            account_name=event.data["account"],
            folder_type="DOCUMENTS",
        )

    async def _on_email_added(self, event: Event):
        await self.process_document(
            file_path=Path(event.data["path"]),
            account_name=event.data["account"],
            folder_type="EMAILS",
        )

    async def _on_rfi_file_added(self, event: Event):
        """User dropped an RFI source file — extract intel, then queue fill_rfi task."""
        file_path = Path(event.data["path"])
        account_name = event.data["account"]
        await self.process_document(
            file_path=file_path,
            account_name=account_name,
            folder_type="RFI",
            doc_type_override="rfi",
        )
        # After extraction, queue the RFI fill task (generates the filled response copy)
        self.event_bus.publish(Event(
            type="rfi.source.ready",
            source="brain.documents",
            data={"account": account_name, "path": str(file_path),
                  "filename": file_path.name}
        ))

    async def _on_presales_file_added(self, event: Event):
        """User dropped a file into a presales section — extract and cascade downstream."""
        file_path = Path(event.data["path"])
        account_name = event.data["account"]
        section = event.data.get("section", "")
        await self.process_document(
            file_path=file_path,
            account_name=account_name,
            folder_type=section,
        )
        # Trigger relevant downstream section refreshes for this section
        cascade_tasks = SECTION_CASCADE_MAP.get(section, [])
        for task_type in cascade_tasks:
            self.event_bus.publish(Event(
                type="presales.cascade.requested",
                source="brain.documents",
                data={"account": account_name, "task": task_type, "trigger": section}
            ))

    async def process_document(
        self, file_path: Path, account_name: str,
        folder_type: str = "DOCUMENTS",
        doc_type_override: Optional[str] = None,
    ):
        """Full pipeline: read → classify → NLP extract (adaptive model) → update account files.

        Model routing:
          - Audio (.mp3/.wav/.m4a): Whisper large-v3-turbo (transcription)
          - Video (.mp4/.mov etc):  Cosmos Reason2 8B (frame analysis)
          - Short text (< 8K chars): Step 3.5 Flash (fast reasoning)
          - Long text (≥ 8K chars):  Nemotron 3 Super 120B (1M-token context, no truncation)
          - Code/config files:       Qwen2.5 Coder 32B
        """
        if not file_path.exists():
            return

        ext = file_path.suffix.lower()
        all_supported = SUPPORTED_EXTS | AUDIO_EXTS | VIDEO_EXTS
        if ext not in all_supported:
            self.logger.debug("Unsupported file type, skipping", path=str(file_path))
            return

        self.logger.info("Processing document", file=file_path.name,
                         account=account_name, folder=folder_type)

        try:
            # Step 1: Extract text (or transcribe audio/video)
            text, task_type = await self._extract_text_adaptive(file_path)
            if not text or len(text.strip()) < 50:
                self.logger.warning("Document too short or empty", file=file_path.name)
                return

            # Step 2: Classify document type (override for known cases like RFP)
            doc_type = doc_type_override or self._classify_document(file_path.name, text)

            # Step 3: NLP extraction — pick model based on content length
            if task_type in ("audio", "video"):
                # Already transcribed; use Step 3.5 Flash for extraction
                extraction_task = "reasoning"
            elif len(text) >= LONG_CONTENT_THRESHOLD:
                extraction_task = "long_context"  # Nemotron 120B for large docs
            else:
                extraction_task = "text"           # Step 3.5 Flash for shorter docs

            extracted = await self._extract_intelligence(
                text, doc_type, account_name, extraction_task
            )

            # Step 4: Write intelligence to account folder
            account_dir = Path(self.config.workspace_root) / "ACCOUNTS" / account_name
            await self._update_account(account_dir, account_name, file_path, doc_type, extracted)

            # Step 5: Trigger downstream skills based on doc type
            self._trigger_skills(account_name, doc_type, extracted)

            self.logger.info("Document processed", file=file_path.name,
                             account=account_name, doc_type=doc_type,
                             model_used=extraction_task)

        except Exception as e:
            self.logger.error("Document processing failed",
                              file=file_path.name, error=str(e))

    # ------------------------------------------------------------------
    # Text extraction (adaptive — returns text + task_type)
    # ------------------------------------------------------------------

    async def _extract_text_adaptive(self, file_path: Path) -> tuple:
        """Extract text from file. Returns (text, task_type) tuple.

        task_type indicates which LLM model class should be used:
          "text"        — short text, Step 3.5 Flash
          "long_context"— long text, Nemotron 120B
          "audio"       — transcribed audio
          "video"       — video frame analysis
          "code"        — code/config files
        """
        ext = file_path.suffix.lower()

        if ext in AUDIO_EXTS:
            text = await self._transcribe_audio(file_path)
            return text, "audio"

        if ext in VIDEO_EXTS:
            # For video we pass the path to Cosmos — return path as placeholder
            return f"[VIDEO FILE: {file_path.name}]", "video"

        if ext in (".py", ".js", ".ts", ".json", ".yaml", ".yml", ".sh", ".sql"):
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            return text, "code"

        if ext in (".txt", ".md"):
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            return text, "text"

        if ext == ".docx":
            text = self._extract_docx(file_path)
            return text, "text"

        if ext == ".pdf":
            text = self._extract_pdf(file_path)
            return text, "text"

        if ext == ".pptx":
            text = self._extract_pptx(file_path)
            return text, "text"

        if ext == ".xlsx":
            text = self._extract_xlsx(file_path)
            return text, "text"

        return "", "text"

    async def _transcribe_audio(self, file_path: Path) -> str:
        """Transcribe audio via NVIDIA Whisper large-v3-turbo."""
        if not self._llm_client:
            self._llm_client = self._build_llm_client()
        if not self._llm_client:
            return ""
        # Pass filename as content — Whisper endpoint handles binary in production
        # For now, return a placeholder so downstream extraction still fires
        self.logger.info("Audio transcription queued", file=file_path.name)
        return f"[AUDIO TRANSCRIPT PENDING: {file_path.name}]"

    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from .pptx slides."""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            self.logger.warning("python-pptx not installed — run: pip install python-pptx")
            return ""
        except Exception as e:
            self.logger.warning("PPTX extraction failed", error=str(e))
            return ""

    def _extract_xlsx(self, file_path: Path) -> str:
        """Extract text from .xlsx spreadsheet."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            texts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            return "\n".join(texts)
        except ImportError:
            self.logger.warning("openpyxl not installed — run: pip install openpyxl")
            return ""
        except Exception as e:
            self.logger.warning("XLSX extraction failed", error=str(e))
            return ""

    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from .docx without external deps."""
        try:
            with zipfile.ZipFile(file_path, "r") as z:
                with z.open("word/document.xml") as f:
                    tree = ET.parse(f)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            texts = [node.text for node in tree.findall(".//w:t", ns) if node.text]
            return " ".join(texts)
        except Exception as e:
            self.logger.warning("DOCX extraction failed", error=str(e))
            return ""

    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF — tries pypdf, falls back to pdfminer."""
        try:
            import pypdf
            reader = pypdf.PdfReader(str(file_path))
            return "\n".join(
                page.extract_text() or "" for page in reader.pages
            )
        except ImportError:
            pass
        try:
            from pdfminer.high_level import extract_text as pdfminer_extract
            return pdfminer_extract(str(file_path))
        except ImportError:
            pass
        self.logger.warning("No PDF library available — install pypdf or pdfminer.six")
        return ""

    # ------------------------------------------------------------------
    # Classification
    # ------------------------------------------------------------------

    def _classify_document(self, filename: str, text: str) -> str:
        fname = filename.lower()
        text_lower = text[:500].lower()

        if any(k in fname for k in ("rfp", "request_for_proposal", "rfi", "rfq")):
            return "rfp"
        if any(k in fname for k in ("contract", "agreement", "nda", "msa", "sow")):
            return "contract"
        if any(k in fname for k in ("brief", "overview", "background")):
            return "brief"
        if any(k in fname for k in ("email", "thread", "mail")):
            return "email"
        if "request for proposal" in text_lower or "rfp" in text_lower:
            return "rfp"
        if "agreement" in text_lower and "parties" in text_lower:
            return "contract"
        return "general"

    # ------------------------------------------------------------------
    # NLP extraction
    # ------------------------------------------------------------------

    async def _extract_intelligence(
        self, text: str, doc_type: str, account_name: str,
        task_type: str = "text"
    ) -> Dict:
        """Use NVIDIA LLM to extract structured intelligence from document text.

        task_type controls which model is used:
          "text"         — Step 3.5 Flash (short docs, default)
          "long_context" — Nemotron 120B (docs ≥ 8K chars, no truncation)
          "reasoning"    — Step 3.5 Flash (for audio transcripts)
          "code"         — Qwen2.5 Coder (code/config files)
        """
        if not self._llm_client:
            self._llm_client = self._build_llm_client()

        if not self._llm_client:
            return {}

        # For long_context model: pass full text (1M token window handles it)
        # For other models: cap at 12K chars to stay within context
        if task_type == "long_context":
            doc_text = text  # full document — Nemotron 120B handles 1M tokens
        else:
            doc_text = text[:12000]

        prompt = (
            f"Account: {account_name}\n"
            f"Document type: {doc_type}\n\n"
            f"{DOCUMENT_TYPE_PROMPTS.get(doc_type, DOCUMENT_TYPE_PROMPTS['general'])}\n\n"
            f"Document text:\n{doc_text}"
        )

        try:
            response = await self._llm_client.generate_with_routing(
                messages=[{"role": "user", "content": prompt}],
                task_type=task_type,
                source="background",
            )
            import re as _re
            json_match = _re.search(r'\{.*\}', response, _re.DOTALL)
            if json_match:
                return json.loads(json_match.group(0))
        except Exception as e:
            self.logger.warning("LLM extraction failed", task_type=task_type, error=str(e))

        return {}

    def _build_llm_client(self):
        try:
            from jarvis.llm.llm_client import LLMClient
            return LLMClient(self.config)
        except Exception:
            return None

    # ------------------------------------------------------------------
    # Account updates
    # ------------------------------------------------------------------

    async def _update_account(
        self, account_dir: Path, account_name: str,
        file_path: Path, doc_type: str, extracted: Dict
    ):
        date_str = datetime.now().strftime("%Y-%m-%d")
        intel_dir = account_dir / "INTEL"
        intel_dir.mkdir(exist_ok=True)

        # Write extracted intelligence as a markdown file
        out_name = f"{date_str}_{file_path.stem}_{doc_type}_intel.md"
        out_file = intel_dir / out_name
        out_file.write_text(
            f"# Document Intelligence: {file_path.name}\n"
            f"**Account:** {account_name}  \n"
            f"**Type:** {doc_type}  \n"
            f"**Processed:** {date_str}  \n\n"
            f"```json\n{json.dumps(extracted, indent=2)}\n```\n",
            encoding="utf-8"
        )

        # Update actions.md with extracted action items
        action_items = extracted.get("action_items", [])
        if action_items:
            actions_file = account_dir / "actions.md"
            with open(actions_file, "a", encoding="utf-8") as f:
                for item in action_items:
                    f.write(f"- [ ] [{date_str}] {item}\n")

        # Update contacts.json
        people = extracted.get("stakeholders", extracted.get("people", []))
        if people:
            contacts_file = account_dir / "contacts.json"
            try:
                data = json.loads(contacts_file.read_text()) if contacts_file.exists() \
                    else {"account": account_name, "contacts": []}
                existing_names = {c.get("name", "").lower() for c in data["contacts"]}
                for person in people:
                    name = person if isinstance(person, str) else person.get("name", "")
                    if name and name.lower() not in existing_names:
                        data["contacts"].append({
                            "name": name,
                            "role": person.get("role", "") if isinstance(person, dict) else "",
                            "source": file_path.name,
                            "first_seen": date_str,
                        })
                        existing_names.add(name.lower())
                contacts_file.write_text(json.dumps(data, indent=2))
            except Exception as e:
                self.logger.warning("Contacts update failed", error=str(e))

        # Log activity
        activities_file = account_dir / "activities.jsonl"
        with open(activities_file, "a", encoding="utf-8") as f:
            f.write(json.dumps({
                "timestamp": datetime.now().isoformat(),
                "source": "document_processor",
                "doc_type": doc_type,
                "file": file_path.name,
                "intel_file": out_name,
            }) + "\n")

    # ------------------------------------------------------------------
    # Downstream skill triggers
    # ------------------------------------------------------------------

    def _trigger_skills(self, account_name: str, doc_type: str, extracted: Dict):
        """Fire skill-trigger events based on what was found in the document."""
        if doc_type == "rfp":
            self.event_bus.publish(Event(
                type="skill.trigger.request",
                source="brain.documents",
                data={"skill": "proposal_generator", "account": account_name,
                      "reason": "RFP detected in DOCUMENTS/"}
            ))
        if extracted.get("competitors"):
            self.event_bus.publish(Event(
                type="skill.trigger.request",
                source="brain.documents",
                data={"skill": "battlecards", "account": account_name,
                      "competitors": extracted["competitors"],
                      "reason": "Competitors mentioned in document"}
            ))
        self.event_bus.publish(Event(
            type="document.processed",
            source="brain.documents",
            data={"account": account_name, "doc_type": doc_type}
        ))
